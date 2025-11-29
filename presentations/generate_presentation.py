#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_presentation.py
Generates a cleaned PPTX presentation 'Carbon_GDP_Presentation.pptx' based on the user's slides.
Requires: python-pptx
Install: pip install python-pptx
Place image files named: image1.png, image2.png, image3.png in the same folder before running.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

OUTFILE = "Carbon_GDP_Presentation.pptx"

prs = Presentation()
# Set slide size to standard 16:9 (optional)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# helper to add title slide
def add_title_slide(title, subtitle, img_path=None, speaker_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Title box
    left = Inches(0.6)
    top = Inches(0.5)
    width = Inches(7.5)
    height = Inches(1.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x0B, 0x3A, 0x66)  # dark blue

    # subtitle
    sub_box = slide.shapes.add_textbox(left, Inches(2.0), width, Inches(0.6))
    st = sub_box.text_frame
    st.text = subtitle
    st.paragraphs[0].font.size = Pt(18)
    st.paragraphs[0].font.name = 'Arial'
    st.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x66, 0x99)

    # right-side image if provided
    if img_path:
        try:
            slide.shapes.add_picture(img_path, Inches(9.0), Inches(0.4), width=Inches(3.1), height=Inches(3.6))
        except Exception as e:
            print("Warning: can't load image:", img_path, e)

    if speaker_note:
        slide.notes_slide.notes_text_frame.text = speaker_note
    return slide

# helper to add bullet slide
def add_bullet_slide(title, bullets, img_path=None, speaker_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.5), Inches(0.6)).text_frame
    t.text = title
    t.paragraphs[0].font.size = Pt(28)
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.name = 'Arial'
    # bullets box
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(7.8), Inches(5.5)).text_frame
    tx.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tx.paragraphs[0]
            p.text = b
        else:
            p = tx.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    # image on right if provided
    if img_path:
        try:
            slide.shapes.add_picture(img_path, Inches(9.0), Inches(1.0), width=Inches(3.1), height=Inches(4.5))
        except Exception as e:
            print("Warning: can't load image:", img_path, e)
    if speaker_note:
        slide.notes_slide.notes_text_frame.text = speaker_note
    return slide

# 1. Title
add_title_slide(
    title="碳排放驱动与 GDP 脱钩路径",
    subtitle="（汇报人：张成浩） — 简化版，便于宣讲",
    img_path="image1.png",
    speaker_note="开场：交代研究背景、目的与逻辑。今天重点讲清楚“为什么重要”和“我们能做什么”。"
)

# 2. 目录
add_bullet_slide(
    title="目录",
    bullets=[
        "01 背景与核心问题",
        "02 理论工具（Kaya / LMDI / EKC）",
        "03 脱钩检验与主要发现",
        "04 政策建议与实践路径",
        "05 结论、AI 与专业反思"
    ],
    img_path="image2.png",
    speaker_note="快速导航：告诉听众节奏与时间分配。"
)

# 3. 背景与问题（精简）
add_bullet_slide(
    title="背景与核心问题",
    bullets=[
        "全球与中国提出碳达峰、碳中和目标（2030/2060）",
        "问题：GDP 增长与碳排放能否脱钩？哪些因素驱动碳排放？",
        "我们关注：可操作的量化工具与政策路径"
    ],
    img_path=None,
    speaker_note="把长段文字压缩成 3 条要点，便于听众跟进。"
)

# 4. 理论工具：Kaya 简述
add_bullet_slide(
    title="理论工具：Kaya 恒等式（简化）",
    bullets=[
        "CO2 = 人口 × 人均 GDP × 能源强度 × 单位能耗碳排放系数",
        "Kaya 把复杂问题分解为可量化的四个驱动项，方便政策针对性设计",
        "演讲提示：举例说明每一项如何被政策或技术影响"
    ],
    img_path="image3.png",
    speaker_note="讲解公式并配合一个具体例子（如提高能效、替代能源减少 CO2/能耗系数）。"
)

# 5. 方法快速说明（LMDI / EKC）
add_bullet_slide(
    title="方法快速说明：LMDI 与 EKC",
    bullets=[
        "LMDI：将排放变化分解为各驱动因素的贡献（无残差、可累加）",
        "EKC：环境库兹涅茨曲线，讨论经济增长与环境压力的关系形态",
        "本报告以 LMDI 分解为主，结合 EKC 的趋势解释"
    ],
    img_path=None,
    speaker_note="补充：LMDI 的优点是可量化每一项贡献，便于政策效果评估。"
)

# 6. 主要发现（示例化）
add_bullet_slide(
    title="脱钩检验：主要发现（示例）",
    bullets=[
        "总体上存在“相对脱钩”趋势，但地区与行业差异显著",
        "能效提升与能源结构改善为主要减排驱动力",
        "短期内 GDP 快速增长仍可能伴随排放增长，需结构性政策配合"
    ],
    img_path=None,
    speaker_note="此处把原始数据表格结论化，现场展示时只强调要点与一两组数据即可。"
)

# 7. 区域与时段异质性（精炼）
add_bullet_slide(
    title="区域与时段异质性",
    bullets=[
        "东部：能源强度下降快，脱钩进程相对领先",
        "中西部：产业结构调整与能源替代潜力大，但起点高",
        "十三五后：政策与技术推动结构性改善"
    ],
    img_path=None,
    speaker_note="如果需要，把某一张原始图表截图放在备注或现场展示板上。"
)

# 8. 政策建议（条目化）
add_bullet_slide(
    title="政策建议（精炼）",
    bullets=[
        "优化能源结构：加快清洁能源替代，重点在电力与重工业",
        "提升能效：推广技术改造与过程能耗管理",
        "区域协同：通过财政/市场手段引导产业转型",
        "公众与企业激励：绿色金融、碳市场、节能激励措施"
    ],
    img_path=None,
    speaker_note="每条建议配 1-2 个可落地的小措施（如碳预算试点、工业能效补贴）。"
)

# 9. 面向节能减排的行动清单（可用于 Q&A）
add_bullet_slide(
    title="面向节能减排，我们需要做（行动清单）",
    bullets=[
        "优化供给侧：发展低碳能源，减少化石能源依赖",
        "推动高质量发展：产业升级，降低能源密集型产出比重",
        "提升能效：建筑、交通、工业重点领域节能改造",
        "倡导绿色生活：节能习惯与消费选择"
    ],
    img_path=None,
    speaker_note="把公众能做的事情列出，便于演讲最后互动问答使用。"
)

# 10. AI 与专业/生活的辩证思考（短小）
add_bullet_slide(
    title="AI 与专业、生活的辩证思考（简短）",
    bullets=[
        "发现：AI 提供强大分析/预测与自动化能力，但也可能增加能耗与伦理问题",
        "小建议（问题解决框架）：Identify -> Analyze -> Propose -> Monitor（示例：用 AI 优化电网调度、但同时监测算力能耗）",
        "提醒：技术是工具，政策与价值引导决定技术带来的社会后果"
    ],
    img_path=None,
    speaker_note="这一页讲两分钟：既肯定 AI 的价值，也提示它带来的新问题，并给出简单的落地思路。"
)

# 11. 结论与致谢
add_bullet_slide(
    title="结论与致谢",
    bullets=[
        "结论：实现 GDP 与碳排放的稳定脱钩需要技术、结构与政策三管齐下",
        "建议重点：清洁能源、能效升级、区域协同与激励机制",
        "感谢聆听，欢迎提问"
    ],
    img_path="image1.png",
    speaker_note="结束语：呼应开头目标，打开 Q&A。"
)

# Save
prs.save(OUTFILE)
print("Saved presentation to", OUTFILE)